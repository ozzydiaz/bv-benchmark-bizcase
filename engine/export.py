"""
engine/export.py
================
Export the business case to PowerPoint (.pptx) and pre-filled Excel (.xlsx).

PowerPoint output
-----------------
Two or three slides:
  Slide 1 — Executive Summary (KPI cards + dual 5Y/10Y stacked-bar+line chart)
  Slide 2 — Detailed Financial Case (annual cashflow table + cumulative savings)
  Slide 3 — Scenario Comparison (only when scenarios=[…] is supplied)

Excel output
------------
Writes the engine's computed inputs back into a copy of the BV Benchmark
Business Case template, filling every yellow input cell so the workbook
produces the same results as the Python engine when opened and saved.
When scenarios=[…] is supplied a `Scenario_Comparison` sheet is appended
with one column per scenario and headline KPIs as rows.

Usage
-----
>>> from engine.export import build_pptx, build_excel
>>> pptx_bytes = build_pptx(summary, fc, inputs, benchmarks, scenarios=alts)
>>> xlsx_bytes = build_excel(inputs, benchmarks, scenarios=alts)

A scenario is a dict with at least:
  {"label": str, "summary": BusinessCaseSummary,
   "aco": list[float], "ecif": list[float], "num_dc_exit": int}
Base scenario is rendered automatically from the positional `summary`.
"""
from __future__ import annotations

import io
import math
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from engine.models import BenchmarkConfig, BusinessCaseInputs
    from engine.financial_case import FinancialCase
    from engine.outputs import BusinessCaseSummary

# Template workbook — path relative to this file so it works on Streamlit Cloud
_TEMPLATE_PATH = Path(__file__).parent.parent / "Template_BV Benchmark Business Case v6.xlsm"

# ---------------------------------------------------------------------------
# Colour palette (consistent with results.py)
# ---------------------------------------------------------------------------
_C_DARK_BG   = "1F1F1F"
_C_PANEL     = "2D2D2D"
_C_SQ        = "FF6B35"   # on-prem coral
_C_AZ        = "0078D4"   # azure blue
_C_CAPEX     = "4A4A6A"   # slate
_C_OPEX      = "C76C00"   # burnt orange
_C_AZURE_BAR = "50B0F0"   # sky blue
_C_MIGRATION = "FFC107"   # amber
_C_WHITE     = "FFFFFF"
_C_GREY      = "AAAAAA"
_C_GREEN     = "00B050"
_C_AMBER     = "FFC000"
_C_RED       = "FF0000"


# ---------------------------------------------------------------------------
# PowerPoint helpers
# ---------------------------------------------------------------------------

def _rgb(hex6: str):
    from pptx.util import Pt  # noqa: F401 – ensure pptx importable
    from pptx.dml.color import RGBColor
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _inches(n: float):
    from pptx.util import Inches
    return Inches(n)


def _pt(n: float):
    from pptx.util import Pt
    return Pt(n)


def _add_textbox(slide, left, top, width, height, text, font_size=11,
                 bold=False, color=_C_WHITE, align="left"):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return txb


def _add_rect(slide, left, top, width, height, fill_hex: str, line_hex: str | None = None):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: F401
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
        shape.line.width = _pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _kpi_card(slide, left: float, top: float, label: str, value: str,
              sub: str = "", accent: str = _C_AZ):
    """Render a KPI card: coloured top accent bar + label + value + sub."""
    card_w, card_h = 1.55, 1.0
    # accent bar
    _add_rect(slide, left, top, card_w, 0.06, accent)
    # card body
    _add_rect(slide, left, top + 0.06, card_w, card_h - 0.06, _C_PANEL)
    # label
    _add_textbox(slide, left + 0.05, top + 0.10, card_w - 0.1, 0.22,
                 label, font_size=7, color=_C_GREY)
    # value
    _add_textbox(slide, left + 0.05, top + 0.30, card_w - 0.1, 0.40,
                 value, font_size=16, bold=True, color=_C_WHITE)
    if sub:
        _add_textbox(slide, left + 0.05, top + 0.72, card_w - 0.1, 0.22,
                     sub, font_size=7, color=_C_GREY)


def _bar_chart_image(summary, horizon: int, width_px=900, height_px=380) -> bytes:
    """
    Render the stacked bar + line chart as a PNG image using plotly's static
    image export (kaleido).  Returns PNG bytes.
    """
    try:
        import plotly.graph_objects as go
        import plotly.io as pio

        years = list(range(1, horizon + 1))
        az_capex = summary.az_cf_capex_by_year[1: horizon + 1]
        az_opex  = summary.az_cf_opex_by_year[1: horizon + 1]
        az_azure = summary.az_cf_azure_by_year[1: horizon + 1]
        az_mig   = [v if v != 0 else None for v in summary.az_cf_migration_by_year[1: horizon + 1]]
        sq_total = summary.sq_cf_by_year[1: horizon + 1]
        az_total = summary.az_cf_by_year[1: horizon + 1]

        fig = go.Figure()
        fig.add_trace(go.Bar(name="Retained CAPEX", x=years, y=az_capex,
                             marker_color=f"#{_C_CAPEX}"))
        fig.add_trace(go.Bar(name="Retained OPEX",  x=years, y=az_opex,
                             marker_color=f"#{_C_OPEX}"))
        fig.add_trace(go.Bar(name="Azure Consumption", x=years, y=az_azure,
                             marker_color=f"#{_C_AZURE_BAR}"))
        fig.add_trace(go.Bar(name="Migration", x=years, y=az_mig,
                             marker_color=f"#{_C_MIGRATION}"))
        fig.add_trace(go.Scatter(name="On-Prem (SQ)", x=years, y=sq_total,
                                 mode="lines+markers",
                                 line=dict(color=f"#{_C_SQ}", width=3), marker=dict(size=6)))
        fig.add_trace(go.Scatter(name="Azure Total", x=years, y=az_total,
                                 mode="lines+markers",
                                 line=dict(color=f"#{_C_AZ}", width=3, dash="dot"),
                                 marker=dict(size=6)))
        fig.update_layout(
            barmode="stack",
            paper_bgcolor="#1F1F1F",
            plot_bgcolor="#1F1F1F",
            font=dict(color="#FFFFFF", size=11),
            xaxis=dict(title="Year", gridcolor="#444444", tickmode="linear", tick0=1, dtick=1),
            yaxis=dict(title="USD", gridcolor="#444444", tickformat="$,.0f"),
            legend=dict(orientation="h", y=-0.25, font=dict(size=9)),
            margin=dict(l=60, r=20, t=20, b=80),
        )
        return pio.to_image(fig, format="png", width=width_px, height=height_px)
    except Exception:
        return b""  # kaleido not installed — skip chart image


# ---------------------------------------------------------------------------
# PowerPoint builder
# ---------------------------------------------------------------------------

def build_pptx(
    summary: "BusinessCaseSummary",
    fc: "FinancialCase",
    inputs: "BusinessCaseInputs",
    benchmarks: "BenchmarkConfig",
    scenarios: "list[dict] | None" = None,
) -> bytes:
    """
    Build a PowerPoint deck and return as bytes.

    Slide 1 — Executive Summary
      • KPI cards: CF NPV 10Y / CF NPV 5Y / P&L NPV 10Y / ROI 10Y / Payback / Yr10 Savings
      • Stacked bar + line chart: 5-year and 10-year horizons (side-by-side)
      • Cumulative savings headline

    Slide 2 — Annual Cash Flow Detail
      • Annual table (Y1–Y10): SQ CF | Azure CF | Annual Savings | Cumulative Savings
      • Waterfall breakdown summary
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    client = inputs.engagement.client_name
    pb_str = f"{summary.payback_years:.1f} yrs" if summary.payback_years else "N/A"
    pb_accent = _C_GREEN if (summary.payback_years or 99) <= 3 else (_C_AMBER if (summary.payback_years or 99) <= 5 else _C_RED)
    roi_accent = _C_GREEN if summary.roi_10yr > 0.5 else (_C_AMBER if summary.roi_10yr > 0 else _C_RED)
    npv_accent = _C_GREEN if summary.npv_cf_10yr > 0 else _C_RED

    # ==================================================================
    # SLIDE 1 — Executive Summary
    # ==================================================================
    s1 = prs.slides.add_slide(blank_layout)

    # Background
    _add_rect(s1, 0, 0, 13.33, 7.5, _C_DARK_BG)

    # Header bar
    _add_rect(s1, 0, 0, 13.33, 0.55, _C_AZ)
    _add_textbox(s1, 0.15, 0.08, 8, 0.40, f"{client} — Azure Migration Business Case", font_size=18, bold=True)
    _add_textbox(s1, 9.0, 0.12, 4.0, 0.30, "Executive Summary", font_size=11, color=_C_WHITE, align="right")

    # KPI cards — row 1
    kpi_top = 0.70
    kpis = [
        ("CF NPV (10-Year)", f"${summary.npv_cf_10yr:,.0f}", "Cash-flow basis", npv_accent),
        ("CF NPV (5-Year)",  f"${summary.npv_cf_5yr:,.0f}",  "Cash-flow basis", npv_accent),
        ("P&L NPV (10-Year)",f"${summary.npv_10yr:,.0f}",   "Depreciation basis", _C_AZ),
        ("10-Year ROI",      f"{summary.roi_10yr:.0%}",     "NPV return on investment", roi_accent),
        ("Payback Period",   pb_str,                         "Cash-flow break-even", pb_accent),
        ("Yr-10 Savings",    f"${summary.savings_yr10:,.0f}", "Annual run-rate", _C_GREEN),
        ("On-Prem Cost/VM",  f"${summary.on_prem_cost_per_vm_yr:,.0f}", "Year 10 / VM / yr", _C_SQ),
        ("Azure Cost/VM",    f"${summary.azure_cost_per_vm_yr:,.0f}",  "Year 10 / VM / yr", _C_AZ),
    ]
    for i, (label, val, sub, accent) in enumerate(kpis):
        col = i % 4
        row = i // 4
        _kpi_card(s1, 0.15 + col * 1.65, kpi_top + row * 1.10, label, val, sub, accent)

    # Chart images
    chart_top = 2.95
    chart_h   = 4.30

    img5  = _bar_chart_image(summary, 5,  width_px=520, height_px=350)
    img10 = _bar_chart_image(summary, 10, width_px=520, height_px=350)

    if img5:
        from pptx.util import Inches as _I
        s1.shapes.add_picture(io.BytesIO(img5),  _I(0.15), _I(chart_top), _I(6.4), _I(chart_h))
    else:
        _add_textbox(s1, 0.15, chart_top + 1.5, 6.4, 0.5,
                     "[5-Year chart — install kaleido for image export]",
                     font_size=9, color=_C_GREY, align="center")

    if img10:
        from pptx.util import Inches as _I
        s1.shapes.add_picture(io.BytesIO(img10), _I(6.78), _I(chart_top), _I(6.40), _I(chart_h))
    else:
        _add_textbox(s1, 6.78, chart_top + 1.5, 6.4, 0.5,
                     "[10-Year chart — install kaleido for image export]",
                     font_size=9, color=_C_GREY, align="center")

    # Chart labels
    _add_textbox(s1, 0.15, chart_top - 0.25, 6.4, 0.25,
                 "5-Year Cost Comparison", font_size=10, bold=True, color=_C_WHITE)
    _add_textbox(s1, 6.78, chart_top - 0.25, 6.4, 0.25,
                 "10-Year Cost Comparison", font_size=10, bold=True, color=_C_WHITE)

    # Legend note
    legend = ("■ Retained CAPEX  ■ Retained OPEX  ■ Azure Consumption  ■ Migration (amber)  "
              "— On-Prem SQ (line)  --- Azure Total (dotted)")
    _add_textbox(s1, 0.15, 7.22, 13.0, 0.25, legend, font_size=7, color=_C_GREY)

    # ==================================================================
    # SLIDE 2 — Annual Cash Flow Detail
    # ==================================================================
    s2 = prs.slides.add_slide(blank_layout)
    _add_rect(s2, 0, 0, 13.33, 7.5, _C_DARK_BG)
    _add_rect(s2, 0, 0, 13.33, 0.55, _C_AZ)
    _add_textbox(s2, 0.15, 0.08, 8, 0.40, f"{client} — Annual Cash Flow & Savings", font_size=18, bold=True)

    # Table header
    col_x = [0.15, 2.25, 4.35, 6.45, 8.55, 10.65]
    col_w = [1.9, 2.0, 2.0, 2.0, 2.0, 2.0]
    headers = ["Year", "SQ Total CF", "Azure Total CF", "Retained OPEX", "Azure Costs", "Cum. Savings"]
    row_y = 0.65
    _add_rect(s2, 0.1, row_y, 13.0, 0.28, _C_AZ)
    for i, h in enumerate(headers):
        _add_textbox(s2, col_x[i], row_y + 0.03, col_w[i], 0.24, h, font_size=8, bold=True, align="right" if i > 0 else "left")

    cumulative = 0.0
    for yr in range(1, 11):
        row_y += 0.30
        bg = _C_PANEL if yr % 2 == 0 else "252525"
        _add_rect(s2, 0.1, row_y, 13.0, 0.28, bg)
        cumulative += summary.annual_cf_savings[yr]
        cum_color = _C_GREEN if cumulative >= 0 else _C_RED
        row_vals = [
            (f"Year {yr}", "left", _C_WHITE),
            (f"${summary.sq_cf_by_year[yr]:>12,.0f}", "right", _C_SQ),
            (f"${summary.az_cf_by_year[yr]:>12,.0f}", "right", _C_AZ),
            (f"${summary.az_cf_opex_by_year[yr]:>12,.0f}", "right", _C_GREY),
            (f"${summary.az_cf_azure_by_year[yr]:>12,.0f}", "right", _C_AZURE_BAR),
            (f"${cumulative:>12,.0f}", "right", cum_color),
        ]
        for i, (txt, align, col) in enumerate(row_vals):
            _add_textbox(s2, col_x[i], row_y + 0.03, col_w[i], 0.24, txt,
                         font_size=8, color=col, align=align)

    # Totals row
    row_y += 0.30
    _add_rect(s2, 0.1, row_y, 13.0, 0.28, "003366")
    totals = [
        ("10-Year Total", "left", _C_WHITE),
        (f"${summary.total_sq_cf_10yr:>12,.0f}", "right", _C_SQ),
        (f"${summary.total_az_cf_10yr:>12,.0f}", "right", _C_AZ),
        ("", "right", _C_WHITE),
        ("", "right", _C_WHITE),
        (f"${cumulative:>12,.0f}", "right", _C_GREEN if cumulative >= 0 else _C_RED),
    ]
    for i, (txt, align, col) in enumerate(totals):
        _add_textbox(s2, col_x[i], row_y + 0.03, col_w[i], 0.24, txt,
                     font_size=8, bold=True, color=col, align=align)

    # Waterfall summary
    wf_top = row_y + 0.55
    _add_textbox(s2, 0.15, wf_top, 13.0, 0.25, "Average Annual Savings by Category (P&L, 10-Year)",
                 font_size=9, bold=True)
    wf_top += 0.28
    wf_cols = list(summary.waterfall.items())
    for i, (cat, val) in enumerate(wf_cols):
        col = _C_WHITE
        if "Reduction" in cat:
            col = _C_GREEN
        elif "Increase" in cat or "Azure Case" in cat:
            col = _C_SQ
        x_pos = 0.15 + (i % 4) * 3.25
        y_pos = wf_top + (i // 4) * 0.32
        _add_textbox(s2, x_pos, y_pos, 3.1, 0.28,
                     f"{cat}: ${val:,.0f}/yr", font_size=8, color=col)

    # ==================================================================
    # SLIDE 3 — Scenario Comparison (only when alts supplied)
    # ==================================================================
    if scenarios:
        _add_scenario_slide(prs, blank_layout, summary, scenarios, client)

    # ------------------------------------------------------------------
    # Speaker notes — audit trail for every slide
    # ------------------------------------------------------------------
    pb_str_note = f"{summary.payback_years:.1f} yrs" if summary.payback_years else "Does not break even within 10Y"
    _set_notes(s1,
        f"Source: BV Benchmark Business Case engine, Layer 3 parity = 0.\n"
        f"Customer: {client}.\n"
        f"WACC: {benchmarks.wacc:.1%}; perpetual growth: {benchmarks.perpetual_growth_rate:.1%}.\n"
        f"CF NPV 10Y = sum of discounted (SQ_total_cf − Azure_total_cf) for Y1–Y10.\n"
        f"P&L NPV 10Y includes Gordon Growth terminal value: TV = savings[10] × (1+g) / (wacc − g) discounted to PV.\n"
        f"ROI 10Y = NPV (incl. TV) / NPV of Azure 10Y costs.\n"
        f"Payback: {pb_str_note}.")
    _set_notes(s2,
        "Annual Cash Flow Detail — raw CF series Y1–Y10 used for CF NPV and Payback.\n"
        "SQ Total CF = on-prem retained costs (no migration). Azure Total CF = retained CAPEX + retained OPEX + Azure consumption + migration.\n"
        "Cumulative Savings = running sum of (SQ − Azure). Break-even year = first non-negative cumulative.\n"
        "Waterfall is P&L basis (depreciated CAPEX) averaged over 10Y — not CF.")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _set_notes(slide, text: str) -> None:
    """Attach speaker notes to a slide for export-time audit trail."""
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = text
    except Exception:
        pass  # notes are best-effort


def _add_scenario_slide(prs, blank_layout, base_summary, scenarios: list[dict], client: str) -> None:
    """
    Slide 3 — Scenario comparison grid.

    Renders Base + up to 3 alternative scenarios side-by-side as columns of
    KPI cards. If more than 3 alternates are supplied only the first 3 are
    shown; a footer flags the truncation.
    """
    s = prs.slides.add_slide(blank_layout)
    _add_rect(s, 0, 0, 13.33, 7.5, _C_DARK_BG)
    _add_rect(s, 0, 0, 13.33, 0.55, _C_AZ)
    _add_textbox(s, 0.15, 0.08, 10, 0.40, f"{client} — Scenario Comparison",
                 font_size=18, bold=True)
    _add_textbox(s, 9.0, 0.12, 4.0, 0.30, f"{1 + len(scenarios)} scenarios",
                 font_size=11, color=_C_WHITE, align="right")

    cols: list[tuple[str, object, dict]] = [("Base", base_summary, {})]
    for sc in scenarios[:3]:
        cols.append((sc.get("label", "Alt"), sc.get("summary"), sc))

    n_cols = len(cols)
    col_w = (13.33 - 0.30) / n_cols
    kpi_top = 0.85

    kpi_rows = [
        ("5-Yr CF ROI",     lambda s_: f"{s_.roi_cf:.0%}",                                        "Cash-flow ROI over 5Y horizon"),
        ("Payback (5Y CF)", lambda s_: (f"{s_.payback_cf:.1f} yrs" if s_.payback_cf else ">5 yrs"), "Cash-flow break-even"),
        ("CF NPV (5-Yr)",   lambda s_: f"${s_.npv_cf_5yr:,.0f}",                                  "Cash-flow NPV"),
        ("CF NPV (10-Yr)",  lambda s_: f"${s_.npv_cf_10yr:,.0f}",                                 "Cash-flow NPV"),
        ("P&L NPV (10-Yr)", lambda s_: f"${s_.npv_10yr:,.0f}",                                    "Depreciated basis incl. TV"),
        ("Yr-10 Savings",   lambda s_: f"${s_.savings_yr10:,.0f}",                                "Annual run-rate savings"),
        ("Azure $/VM/yr",   lambda s_: f"${s_.azure_cost_per_vm_yr:,.0f}",                        "Year 10 Azure unit cost"),
    ]
    card_h = 0.78

    for col_i, (label, sm, sc_meta) in enumerate(cols):
        left = 0.15 + col_i * col_w
        accent = _C_AZ if col_i == 0 else _C_GREEN
        _add_rect(s, left, kpi_top, col_w - 0.10, 0.06, accent)
        _add_rect(s, left, kpi_top + 0.06, col_w - 0.10, 0.45, _C_PANEL)
        _add_textbox(s, left + 0.05, kpi_top + 0.10, col_w - 0.20, 0.30,
                     label, font_size=12, bold=True, color=_C_WHITE, align="center")

        if sm is None:
            _add_textbox(s, left + 0.05, kpi_top + 0.55, col_w - 0.20, 0.30,
                         "(no summary)", font_size=9, color=_C_GREY, align="center")
            continue

        for row_i, (rlabel, fn, sub) in enumerate(kpi_rows):
            top = kpi_top + 0.55 + row_i * (card_h + 0.04)
            _add_rect(s, left, top, col_w - 0.10, card_h, _C_PANEL)
            _add_textbox(s, left + 0.05, top + 0.04, col_w - 0.20, 0.20,
                         rlabel, font_size=7, color=_C_GREY)
            try:
                value = fn(sm)
            except Exception:
                value = "—"
            _add_textbox(s, left + 0.05, top + 0.22, col_w - 0.20, 0.34,
                         value, font_size=14, bold=True, color=_C_WHITE)
            _add_textbox(s, left + 0.05, top + 0.55, col_w - 0.20, 0.20,
                         sub, font_size=6, color=_C_GREY)

        if sc_meta:
            aco_total  = sum(abs(x) for x in sc_meta.get("aco", []) or [])
            ecif_total = sum(abs(x) for x in sc_meta.get("ecif", []) or [])
            dc         = sc_meta.get("num_dc_exit", 0)
            footer     = f"DC exit: {dc} · ACO: ${aco_total:,.0f} · ECIF: ${ecif_total:,.0f}"
            _add_textbox(s, left + 0.05, 7.05, col_w - 0.20, 0.25,
                         footer, font_size=7, color=_C_GREY)

    if len(scenarios) > 3:
        _add_textbox(s, 0.15, 7.30, 13.0, 0.18,
                     f"⚠ {len(scenarios) - 3} additional scenario(s) omitted from slide for layout. "
                     f"All {1 + len(scenarios)} scenarios are exported in the Scenario_Comparison sheet of the XLSX.",
                     font_size=7, color=_C_AMBER)

    _set_notes(s,
        f"Scenario grid — {1 + len(scenarios)} scenarios (Base + {len(scenarios)} alt).\n"
        "All KPIs computed by the same Layer-3 engine; only WACC, ACO, ECIF and DC-exit differ across columns.\n"
        "Engine drift vs BA workbook: 0 cells (Customer A baseline).")


# ---------------------------------------------------------------------------
# Excel builder — pre-fill template yellow cells
# ---------------------------------------------------------------------------

def build_excel(
    inputs: "BusinessCaseInputs",
    benchmarks: "BenchmarkConfig",
    template_path: "str | Path | None" = None,
    scenarios: "list[dict] | None" = None,
) -> bytes:
    """
    Write engine inputs into the BV Benchmark workbook template yellow cells
    and return the modified workbook as bytes (.xlsx, macros stripped).

    The resulting file can be opened in Excel to verify or re-run formulas.
    Cells written: 1-Client Variables (D9, D10, D24, D25, D26, D39, D44, D49,
    D54, D66, D67, D68) + 2a-Consumption Plan (D8, D9, D10, E17:N17, E21:N21,
    E22:N22).

    When `scenarios` is supplied, an extra `Scenario_Comparison` sheet is
    appended with one column per scenario and headline KPIs as rows. The
    base inputs are still written to the template's yellow cells so the
    workbook recalculates the **base** scenario when opened in Excel; the
    Scenario_Comparison sheet is values-only (no formulas).
    """
    import openpyxl

    resolved = Path(template_path) if template_path else _TEMPLATE_PATH
    wb = openpyxl.load_workbook(resolved, keep_vba=False, data_only=False)
    cv  = wb["1-Client Variables"]
    cp  = wb["2a-Consumption Plan Wk1"]

    wl  = inputs.workloads[0]         if inputs.workloads         else None
    plan = inputs.consumption_plans[0] if inputs.consumption_plans else None

    # Engagement
    cv["D9"]  = inputs.engagement.client_name
    cv["D10"] = inputs.engagement.local_currency_name

    # Hardware lifecycle
    cv["D24"] = inputs.hardware.depreciation_life_years
    cv["D25"] = inputs.hardware.actual_usage_life_years
    cv["D26"] = inputs.hardware.expected_future_growth_rate
    cv["D27"] = inputs.hardware.hardware_renewal_during_migration_pct

    if wl:
        cv["D39"] = wl.num_vms
        # D40 is BA's manual integer cell; engine carries a fractional topology
        # residual (D42 - num_vms/K11). Round to int for Excel display fidelity.
        cv["D40"] = int(round(wl.num_physical_servers_excl_hosts))
        cv["D44"] = wl.allocated_vcpu
        cv["D49"] = wl.allocated_vmemory_gb
        cv["D54"] = wl.allocated_storage_gb
        cv["D66"] = wl.vcpu_per_core_ratio
        cv["D67"] = wl.pcores_with_windows_server
        cv["D68"] = wl.pcores_with_windows_esu
        if wl.pcores_with_sql_server is not None:
            cv["D70"] = wl.pcores_with_sql_server
        if wl.pcores_with_sql_esu is not None:
            cv["D71"] = wl.pcores_with_sql_esu

    if plan:
        cp["D8"]  = plan.azure_vcpu
        cp["D9"]  = plan.azure_memory_gb
        cp["D10"] = plan.azure_storage_gb

        ramp_cols = ["E", "F", "G", "H", "I", "J", "K", "L", "M", "N"]
        for i, col in enumerate(ramp_cols):
            cp[f"{col}17"] = plan.migration_ramp_pct[i]
            cp[f"{col}21"] = plan.aco_by_year[i]
            cp[f"{col}22"] = plan.ecif_by_year[i]

    # ------------------------------------------------------------------
    # Scenario comparison sheet (values-only)
    # ------------------------------------------------------------------
    if scenarios:
        _write_scenario_sheet(wb, scenarios, benchmarks)

    # Audit-metadata sheet — engine version + drift status
    _write_audit_sheet(wb, inputs, benchmarks, scenarios)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


def _write_scenario_sheet(wb, scenarios: list[dict], benchmarks: "BenchmarkConfig") -> None:
    """
    Append a `Scenario_Comparison` sheet listing all alternate scenarios
    side-by-side. Values only (no formulas) — Excel cannot re-derive these
    from yellow cells because the BA template only models a single scenario.
    """
    sheet_name = "Scenario_Comparison"
    if sheet_name in wb.sheetnames:
        del wb[sheet_name]
    ws = wb.create_sheet(sheet_name)

    # Header
    ws.cell(row=1, column=1, value="KPI / Source")
    ws.cell(row=1, column=2, value="Base (workbook)")
    for col_i, sc in enumerate(scenarios, start=3):
        ws.cell(row=1, column=col_i, value=sc.get("label", f"Scenario {col_i - 2:02d}"))

    rows = [
        ("5-Yr CF ROI",          lambda s_: s_.roi_cf,                "0.0%"),
        ("Payback (5Y CF, yrs)", lambda s_: s_.payback_cf if s_.payback_cf else None, "0.00"),
        ("CF NPV (5-Yr)",        lambda s_: s_.npv_cf_5yr,            "$#,##0"),
        ("CF NPV (10-Yr)",       lambda s_: s_.npv_cf_10yr,           "$#,##0"),
        ("P&L NPV (5-Yr)",       lambda s_: s_.npv_5yr,               "$#,##0"),
        ("P&L NPV (10-Yr)",      lambda s_: s_.npv_10yr,              "$#,##0"),
        ("Yr-10 Annual Savings", lambda s_: s_.savings_yr10,          "$#,##0"),
        ("Azure Cost / VM / yr", lambda s_: s_.azure_cost_per_vm_yr,  "$#,##0"),
        ("On-Prem Cost / VM / yr", lambda s_: s_.on_prem_cost_per_vm_yr, "$#,##0"),
        ("Savings / VM / yr",    lambda s_: s_.savings_per_vm_yr,     "$#,##0"),
        ("Total SQ CF (10Y)",    lambda s_: s_.total_sq_cf_10yr,      "$#,##0"),
        ("Total Az CF (10Y)",    lambda s_: s_.total_az_cf_10yr,      "$#,##0"),
    ]

    # Helper: extract summary; column 2 has no scenario object
    def _summary_for(col_i: int):
        if col_i == 2:
            # Base column is intentionally blank — workbook will recalculate
            # from the yellow cells when opened in Excel.
            return None
        sc = scenarios[col_i - 3]
        return sc.get("summary")

    for row_i, (label, fn, fmt) in enumerate(rows, start=2):
        ws.cell(row=row_i, column=1, value=label)
        for col_i in range(2, 3 + len(scenarios)):
            sm = _summary_for(col_i)
            if sm is None:
                ws.cell(row=row_i, column=col_i, value="(see workbook recalc)")
            else:
                try:
                    val = fn(sm)
                except Exception:
                    val = None
                cell = ws.cell(row=row_i, column=col_i, value=val)
                if val is not None and isinstance(val, (int, float)):
                    cell.number_format = fmt

    # Scenario metadata footer
    meta_row = 2 + len(rows) + 1
    ws.cell(row=meta_row,     column=1, value="— Scenario inputs —").font = ws.cell(row=meta_row, column=1).font.copy(bold=True)
    ws.cell(row=meta_row + 1, column=1, value="WACC")
    ws.cell(row=meta_row + 2, column=1, value="DC exits")
    ws.cell(row=meta_row + 3, column=1, value="ACO total (abs)")
    ws.cell(row=meta_row + 4, column=1, value="ECIF total (abs)")
    ws.cell(row=meta_row + 1, column=2, value=benchmarks.wacc)
    ws.cell(row=meta_row + 1, column=2).number_format = "0.00%"
    for col_i, sc in enumerate(scenarios, start=3):
        # Per-scenario WACC isn't preserved on the dict; show base wacc unless
        # caller stuffs a `wacc` key on the scenario dict.
        ws.cell(row=meta_row + 1, column=col_i, value=sc.get("wacc", benchmarks.wacc)).number_format = "0.00%"
        ws.cell(row=meta_row + 2, column=col_i, value=sc.get("num_dc_exit", 0))
        ws.cell(row=meta_row + 3, column=col_i,
                value=sum(abs(x) for x in (sc.get("aco")  or []))).number_format = "$#,##0"
        ws.cell(row=meta_row + 4, column=col_i,
                value=sum(abs(x) for x in (sc.get("ecif") or []))).number_format = "$#,##0"

    # Column widths
    from openpyxl.utils import get_column_letter
    ws.column_dimensions[get_column_letter(1)].width = 26
    for col_i in range(2, 3 + len(scenarios)):
        ws.column_dimensions[get_column_letter(col_i)].width = 22


def _write_audit_sheet(wb, inputs, benchmarks, scenarios) -> None:
    """Hidden audit-trail sheet — engine version, parity status, key benchmarks."""
    sheet_name = "_Audit"
    if sheet_name in wb.sheetnames:
        del wb[sheet_name]
    ws = wb.create_sheet(sheet_name)
    ws.sheet_state = "hidden"

    rows = [
        ("Engine version",                "v1.5.0-layer3-zero-drift"),
        ("Layer 3 parity",                "0 cells drift vs BA workbook (Customer A baseline)"),
        ("WACC (discount rate)",          benchmarks.wacc),
        ("Perpetual growth (Gordon g)",   benchmarks.perpetual_growth_rate),
        ("Inflation rate",                getattr(benchmarks, "inflation_rate", None)),
        ("Customer",                      inputs.engagement.client_name),
        ("Currency",                      inputs.engagement.local_currency_name),
        ("Workloads modelled",            len(inputs.workloads)),
        ("Consumption plans modelled",    len(inputs.consumption_plans)),
        ("Scenarios in this export",      1 + (len(scenarios) if scenarios else 0)),
        ("Notes",
         "TV: Gordon Growth on P&L savings[10] (see version-history.md for v1.6/v1.7 backlog). "
         "Azure pricing: PAYG list, single flat ACD. RI/SP blending deferred to v1.7."),
    ]
    for r, (label, val) in enumerate(rows, start=1):
        ws.cell(row=r, column=1, value=label)
        ws.cell(row=r, column=2, value=val)

    from openpyxl.utils import get_column_letter
    ws.column_dimensions[get_column_letter(1)].width = 32
    ws.column_dimensions[get_column_letter(2)].width = 60
